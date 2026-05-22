from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


ROOT = Path(__file__).resolve().parents[1]
OUTPUT_PATH = ROOT / "docs" / "NT131_Smart_Parking_FreeRTOS_Hardware.pptx"

ACCENT = RGBColor(0x0F, 0x4C, 0x81)
ACCENT_2 = RGBColor(0xF2, 0x7A, 0x24)
TEXT_DARK = RGBColor(0x1F, 0x29, 0x37)
TEXT_MUTED = RGBColor(0x6B, 0x72, 0x80)
BG = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xE5, 0xE7, 0xEB)


slides = [
    {
        "kind": "title",
        "title": "NT131 Smart Parking",
        "subtitle": "FreeRTOS & Hardware Deep Dive\nKiến trúc phần cứng, luồng I/O, realtime Socket.IO và kết quả thực nghiệm",
        "footer": "NT131 Team",
    },
    {
        "title": "Mục tiêu báo cáo",
        "bullets": [
            "Mô tả chi tiết thiết kế phần cứng và phân chia tác vụ FreeRTOS",
            "Giải thích luồng dữ liệu vào/ra: RFID, ANPR, Gate commands",
            "Trình bày kết quả thực nghiệm: latency, CPU, memory",
            "Nêu rõ contract realtime Socket.IO và tích hợp với backend",
        ],
    },
    {
        "title": "Tổng quan hệ thống",
        "bullets": [
            "Thành phần: Backend API & Realtime, Hardware gateway (ESP32), Database (MongoDB)",
            "Mối liên hệ: Hardware ↔ Socket.IO ↔ Backend ↔ DB",
            "Trọng tâm báo cáo: Hardware (ESP32) và FreeRTOS multitasking",
        ],
    },
    {
        "title": "Thiết kế phần cứng (ESP32 gate controller)",
        "bullets": [
            "ESP32 dev board (WiFi)",
            "MFRC522 RFID reader",
            "Servo motor (barrier)",
            "Camera/ANPR (optional, external)",
            "16x2 LCD (debug)",
            "Giao tiếp mạng: Socket.IO client kết nối tới backend",
            "Yêu cầu: low-latency command handling, non-blocking servo, reliable RFID polling",
        ],
    },
    {
        "title": "Ứng dụng FreeRTOS trong xử lý song song",
        "bullets": [
            "Khai thác kiến trúc lõi kép của ESP32 thông qua hàm xTaskCreatePinnedToCore() để ghim các task lên hai lõi và phân tầng ưu tiên.",
            "Thực tế triển khai gồm 5 task chính (không chỉ 3):",
            "taskRfidPolling (Core 0 - Priority 4): quét MFRC522 chu kỳ nhanh (~50ms) để không bỏ sót thẻ.",
            "taskSocketIO (Core 1 - Priority 3): duy trì Socket.IO, gửi/nhận event, quản lý room/ack.",
            "taskWifiManager (Core 1 - Priority 2): quản lý WiFi, reconnect và trạng thái mạng.",
            "taskServoControl (Core 0 - Priority 2): thực hiện sweep servo theo command, non-blocking (vTaskDelay).",
            "taskLcdDisplay (Core 0 - Priority 1): cập nhật màn hình LCD và hiển thị debug info.",
            "Đồng bộ giữa các task dùng Queue, Semaphore và flags; servo chỉ chạy khi có semaphore/flag.",
        ],
    },
    {
        "title": "Vai trò từng luồng",
        "bullets": [
            "taskRfidPolling: ưu tiên cao nhất, đảm bảo không bỏ lỡ sự kiện quét thẻ.",
            "taskSocketIO: giữ kết nối realtime ổn định và xử lý event từ backend.",
            "taskWifiManager: đảm bảo mạng ổn định, tự reconnect khi mất kết nối.",
            "taskServoControl: thực hiện hành động chấp hành (servo) theo command, sử dụng non-blocking sweep.",
            "taskLcdDisplay: cung cấp thông tin debug/hiển thị trạng thái, ưu tiên thấp nhất.",
            "Tổ chức này tách rõ cảm biến, mạng và chấp hành để hệ thống phản ứng ổn định hơn.",
        ],
    },
    {
        "title": "Thiết kế mạch và kết nối ngoại vi",
        "bullets": [
            "RFID reader MFRC522 kết nối qua SPI để đọc UID nhanh và ổn định",
            "Servo barrier điều khiển bằng PWM, cần nguồn riêng và nối mass chung với ESP32",
            "LCD 16x2 dùng I2C để hiển thị trạng thái hệ thống và debug",
            "WiFi là kênh truyền realtime, ESP32 đóng vai trò gateway",
            "Lưu ý phần cứng: chống nhiễu, kiểm tra nguồn cấp, và tránh xung đột chân I/O",
        ],
    },
    {
        "title": "Luồng sự kiện và đồng bộ trạng thái",
        "bullets": [
            "Sự kiện từ RFID đi qua pipeline: detect → verify → quyết định mở/khóa barrier",
            "Trạng thái hệ thống chuyển theo chuỗi: idle → waiting_scan → approved/rejected → gate_open → gate_closed",
            "Khi backend phản hồi, hardware cập nhật trạng thái và phát tiếp sự kiện realtime",
            "Cách này giúp đảm bảo hệ thống phản ứng nhất quán giữa cảm biến, servo và backend",
        ],
    },
    {
        "title": "Servo control: non-blocking pattern",
        "bullets": [
            "TaskHardware chỉ chạy servo khi nhận được tín hiệu từ cờ trạng thái hoặc Semaphore",
            "Mỗi bước sweep dùng vTaskDelay(pdMS_TO_TICKS(...)) để nhường CPU cho TaskRFID và TaskSocket",
            "Kết quả: không bỏ lỡ quét RFID trong khi barrier hoạt động",
        ],
    },
    {
        "title": "RFID polling & event pipeline",
        "bullets": [
            "taskRfidPolling poll ~100ms; khi detect, push vào queueRfidEvent",
            "taskSocketIO đọc queue, compose canonical envelope và emit tới backend",
            "Envelope chứa eventId, eventName, occurredAt, correlationId, payload",
            "Backend dùng eventId làm idempotency key",
        ],
    },
    {
        "title": "Realtime integration (Socket.IO) — Backend side",
        "bullets": [
            "Socket endpoint: /socket.io (auth via handshake auth.token)",
            "Rooms: operator, simulator (hardware joins as gateway)",
            "Outbound channels: unified realtime.event (envelope) và named events",
            "Inbound hardware commands: backend emits gate.command.sent → hardware xử lý và ack bằng gate.state.changed",
        ],
    },
    {
        "title": "Realtime: Contract và thực thi idempotency",
        "bullets": [
            "Envelope fields required; events cốt lõi: gate.command.sent, gate.state.changed, session.*, rfid.scan.*, slot.*",
            "Quy tắc xử lý trên hardware: lưu lastHandledEventId để bỏ qua duplicate",
            "Ack explicit (result: ack/nack/timeout)",
            "Emit state transitions sau khi cập nhật trạng thái local",
        ],
    },
    {
        "title": "Luồng nghiệp vụ chi tiết: Entry happy path",
        "bullets": [
            "taskRfidPolling phát hiện UID → queue",
            "taskSocketIO emit rfid.scan.requested → Backend xử lý (auth/verify)",
            "Backend tạo/cập nhật parking_session → emit rfid.scan.accepted hoặc rfid.scan.rejected",
            "Nếu accepted: Backend gửi gate.command.sent (open) → hardware nhận và queue servo command",
            "taskServoControl thực hiện mở barrier; khi hoàn thành emit gate.state.changed (open)",
        ],
    },
    {
        "title": "Kết quả thực nghiệm",
        "bullets": [
            "Latency: ~5–10 ms từ backend command tới servo start (giảm so với blocking)",
            "RFID responsiveness: sub-100 ms polling detection",
            "CPU usage: ~20% idle, đủ headroom",
            "Memory: ~8 KB heap cho queues/mutexes, ~18 KB cho stacks",
            "Behavior: non-blocking sweep cho phép tiếp tục poll RFID và nhận event",
        ],
    },
    {
        "title": "Kiến nghị & bước tiếp theo",
        "bullets": [
            "Thêm telemetry cho tasks (stack high-water mark, heap) và gửi về backend",
            "Tối ưu kích thước queue và priority dựa trên tải thực nghiệm",
            "Thêm watchdog task và OTA update secure",
            "Viết test kịch bản latency/regression cho gate commands",
        ],
    },
]


def set_run_font(run, size: int, bold: bool = False, color=TEXT_DARK, name: str = "Aptos"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_background(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.18))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def add_footer(slide, slide_num: int):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(7.02), Inches(12.2), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT
    line.line.fill.background()

    footer = slide.shapes.add_textbox(Inches(0.55), Inches(7.06), Inches(8.5), Inches(0.25))
    p = footer.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "NT131 Smart Parking"
    set_run_font(r, 10, False, TEXT_MUTED)

    num = slide.shapes.add_textbox(Inches(12.0), Inches(7.04), Inches(0.7), Inches(0.25))
    p2 = num.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = str(slide_num)
    set_run_font(r2, 10, False, TEXT_MUTED)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Default slide layout ids: 0 title, 1 title+content
for index, spec in enumerate(slides, start=1):
    if spec.get("kind") == "title":
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_background(slide)

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(7.5), Inches(1.3))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = spec["title"]
        set_run_font(r, 28, True, ACCENT)

        subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.25), Inches(7.5), Inches(1.5))
        stf = subtitle_box.text_frame
        stf.word_wrap = True
        stf.vertical_anchor = MSO_ANCHOR.TOP
        p = stf.paragraphs[0]
        for i, line in enumerate(spec["subtitle"].split("\n")):
            if i:
                p = stf.add_paragraph()
            p.text = line
            if i == 0:
                p.space_after = Pt(8)
            for run in p.runs:
                set_run_font(run, 18, False, TEXT_DARK)

        accent_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.0), Inches(1.1), Inches(3.2), Inches(4.8))
        accent_panel.fill.solid()
        accent_panel.fill.fore_color.rgb = ACCENT
        accent_panel.line.color.rgb = ACCENT

        mini = slide.shapes.add_textbox(Inches(9.35), Inches(1.55), Inches(2.5), Inches(3.8))
        mtf = mini.text_frame
        mtf.word_wrap = True
        p = mtf.paragraphs[0]
        r = p.add_run()
        r.text = "Hardware\nFreeRTOS\nSocket.IO\nRealtime\nExperiments"
        set_run_font(r, 24, True, WHITE)
        p.alignment = PP_ALIGN.CENTER
        mtf.vertical_anchor = MSO_ANCHOR.MIDDLE

        footer = slide.shapes.add_textbox(Inches(0.8), Inches(6.35), Inches(5), Inches(0.3))
        fp = footer.text_frame.paragraphs[0]
        rr = fp.add_run()
        rr.text = spec["footer"]
        set_run_font(rr, 12, True, TEXT_MUTED)
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_background(slide)

        title = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12.0), Inches(0.65))
        tf = title.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"{index} — {spec['title']}"
        set_run_font(r, 24, True, ACCENT)

        content = slide.shapes.add_textbox(Inches(0.9), Inches(1.35), Inches(11.9), Inches(5.3))
        tf = content.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)

        for bullet_index, bullet in enumerate(spec["bullets"]):
            if bullet_index == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.level = 0
            p.space_after = Pt(6)
            p.space_before = Pt(1)
            p.line_spacing = 1.15
            r = p.add_run()
            r.text = bullet
            set_run_font(r, 18 if len(bullet) < 90 else 16, False, TEXT_DARK)
            p.bullet = True

        # Decorative side tag
        tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.45), Inches(6.0), Inches(1.2), Inches(0.42))
        tag.fill.solid()
        tag.fill.fore_color.rgb = ACCENT_2
        tag.line.fill.background()
        tag_text = slide.shapes.add_textbox(Inches(11.45), Inches(6.03), Inches(1.2), Inches(0.2))
        pt = tag_text.text_frame.paragraphs[0]
        pt.alignment = PP_ALIGN.CENTER
        rt = pt.add_run()
        rt.text = "NT131"
        set_run_font(rt, 10, True, WHITE)

    add_footer(slide, index)

prs.core_properties.title = "NT131 Smart Parking - FreeRTOS & Hardware Deep Dive"
prs.core_properties.subject = "Hardware focused slide deck"
prs.core_properties.author = "GitHub Copilot"
prs.core_properties.company = "NT131 Team"
prs.core_properties.comments = "Generated from docs/presentation-slides.md"

OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUTPUT_PATH)
print(f"PPTX created: {OUTPUT_PATH}")
